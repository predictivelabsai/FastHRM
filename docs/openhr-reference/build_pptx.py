"""Generate docs/cbre/cbre_crm_architecture_choices.pptx.

Author: Michael Palys -- Predictive Labs review for CBRE.
Rebuild with:  python docs/build_pptx.py
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
CBRE = DOCS / "cbre"
OUT = CBRE / "cbre_crm_architecture_choices.pptx"
DIAGRAMS = DOCS / "diagrams"

CBRE_GREEN = RGBColor(0x00, 0x3F, 0x2D)
CBRE_ACCENT = RGBColor(0x17, 0xE8, 0x8F)
INK = RGBColor(0x1A, 0x20, 0x2C)
MUTED = RGBColor(0x6C, 0x79, 0x89)
PAPER = RGBColor(0xF5, 0xF7, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDD, 0xE3, 0xEA)


def add_title_slide(prs: Presentation, title: str, subtitle: str, author: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CBRE_GREEN
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.2), Inches(1.5), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CBRE_ACCENT
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sb = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12), Inches(0.7))
    p = sb.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = CBRE_ACCENT

    fb = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12), Inches(0.5))
    p = fb.text_frame.paragraphs[0]
    p.text = f"{author}  |  {date.today().isoformat()}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xCC, 0xD4, 0xDC)


def add_section_header(prs: Presentation, header: str, subhead: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CBRE_GREEN
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = header
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subhead:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(12), Inches(0.7))
        p = sb.text_frame.paragraphs[0]
        p.text = subhead
        p.font.size = Pt(18)
        p.font.color.rgb = CBRE_ACCENT


def add_slide_title(slide, text: str):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CBRE_GREEN
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(1.2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CBRE_ACCENT
    bar.line.fill.background()


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title)
    top = Inches(1.15)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.6))
        p = sb.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(15)
        p.font.italic = True
        p.font.color.rgb = MUTED
        top = Inches(1.75)

    body = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(5.7))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(17)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def add_two_column_slide(prs: Presentation, title: str, left_head: str, left: list[str],
                         right_head: str, right: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title)

    for x, head, bullets, accent in (
        (Inches(0.5), left_head, left, CBRE_GREEN),
        (Inches(6.85), right_head, right, RGBColor(0x22, 0x4A, 0x9D)),
    ):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3),
                                      Inches(6.1), Inches(5.7))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = BORDER
        card.line.width = Pt(0.75)

        ht = slide.shapes.add_textbox(x, Inches(1.45), Inches(6.1), Inches(0.6))
        p = ht.text_frame.paragraphs[0]
        p.text = head
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = accent

        body = slide.shapes.add_textbox(Inches(x.emu / 914400 + 0.05), Inches(2.1),
                                        Inches(5.95), Inches(4.7))
        tf = body.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(14)
            p.font.color.rgb = INK
            p.space_after = Pt(6)


def add_comparison_table(prs: Presentation, title: str, headers: list[str],
                          rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title)
    left, top = Inches(0.5), Inches(1.2)
    width = prs.slide_width - Inches(1.0)
    height = Inches(5.9)
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                        left, top, width, height)
    table = tbl_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CBRE_GREEN
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if r % 2 == 1 else PAPER
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = INK


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title)
    max_w = prs.slide_width - Inches(0.8)
    max_h = prs.slide_height - Inches(1.8)
    pic = slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(1.2), width=max_w)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = int((prs.slide_width - pic.width) / 2)
    pic.top = Inches(1.2)
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.55),
                                      prs.slide_width - Inches(1.0), Inches(0.45))
        p = cb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = MUTED


def add_roadmap_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Sequencing: 12 months, semantic layer + AI agents")
    phases = [
        ("Phase 0", "Wk 0-4",   "Foundations: BFF, event bus, SSO, design system"),
        ("Phase 1", "Wk 4-10",  "Executive dashboard on React Native -- reads SF + legacy"),
        ("Phase 2", "Wk 8-16",  "Frappe CRM pilot -- benchmark against specialist CRE SaaS (e.g. Enaia)"),
        ("Phase 3", "Wk 16-20", "CRE DocTypes (Property, Lease, BrokerEngagement) on Frappe"),
        ("Phase 4", "Wk 18-30", "Accounts mastered + CRE ontology v1 + knowledge graph seeded"),
        ("Phase 5", "Wk 30-42", "AI agents + KG traversal: comp analysis, lease alerts, research"),
        ("Phase 6", "Wk 42-52", "Full workflow move; ontology extended; multi-agent orchestration"),
    ]
    top = Inches(1.3)
    for i, (phase, weeks, desc) in enumerate(phases):
        y = top + Inches(0.78 * i)
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.5), y, Inches(1.4), Inches(0.65))
        chip.fill.solid()
        chip.fill.fore_color.rgb = CBRE_GREEN
        chip.line.fill.background()
        tf = chip.text_frame
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = phase
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        wb = slide.shapes.add_textbox(Inches(2.1), y + Inches(0.08), Inches(1.5), Inches(0.5))
        p = wb.text_frame.paragraphs[0]
        p.text = weeks
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = MUTED

        db = slide.shapes.add_textbox(Inches(3.7), y + Inches(0.08), Inches(9.2), Inches(0.55))
        p = db.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = INK


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        title="CBRE CRM: Architecture Choices",
        subtitle="OSS core + semantic layer + knowledge graph: an AI-ready target architecture.",
        author="Author: Michael Palys",
    )

    # 2. Context
    add_bullet_slide(
        prs,
        "Where CBRE is today",
        [
            "Tableau front-ends layered on SharePoint -- good to see, bad to act.",
            "Proprietary CRM / project tool underneath -- old, isolated, a single point of risk.",
            "Tableau and the proprietary CRM do not talk to each other.",
            "Salesforce gravity is real in the estate -- sales and finance trust it.",
            "Design ambition: React-Native + Material UI based on the new CBRE design system.",
        ],
        subtitle="The as-is: read-only analytics sitting on top of an opaque system of record.",
    )

    # 3. The three needs
    add_bullet_slide(
        prs,
        "The three functional needs",
        [
            'Executive dashboards that let leaders see the business AND kick off workflows -- a "system of action".',
            "Account maintenance -- creation, setup, ownership, lifecycle.",
            "CRM / project management -- pipelines, deals, work packages, timelines.",
        ],
        subtitle="All three asks are real -- they are not the same product.",
    )

    # 4. The two failure modes
    add_two_column_slide(
        prs,
        "Two failure modes on the table",
        left_head="Fail mode A: Let Salesforce rebuild everything",
        left=[
            "Multi-quarter delivery before basics are live.",
            "License + SI + customisation -- seven-figure run rate.",
            "Hard to unwind later. Lock-in tightens while AI rearranges vendor economics.",
            "Data egress, AI SKUs, storage limits -- not your knobs.",
            "Fine for pipeline. Wrong for the whole estate.",
        ],
        right_head="Fail mode B: Full custom composable platform",
        right=[
            "You rebuild accounts, permissions, workflow, reports from zero.",
            "Classic multi-year enterprise programme that never finishes.",
            "80% of what you build is undifferentiated.",
            "Great engineering team evaporates maintaining glue.",
            "Avoid by standing on an OSS substrate (Frappe CRM) -- learn domain features from specialist CRE SaaS like Enaia.",
        ],
    )

    # 5. Why OSS
    add_bullet_slide(
        prs,
        "Why open-source beats a Salesforce-led rebuild",
        [
            'Atlassian lesson: an "extensible platform" you don\'t own can pivot (Server to Cloud) and invalidate years of customisation overnight.',
            "With OSS you own the schema, the upgrade cadence, and the exit cost.",
            "AI value is moving to whoever has the data. Keep the data on your side of the line.",
            "Extensions are code you wrote, not apps you rent -- no marketplace-churn tax.",
            "Exit cost becomes the cost of moving a Postgres database, not a re-platform.",
        ],
        subtitle="Not free in year one. Dramatically cheaper across years two to five.",
    )

    # 6. OSS comparison table
    add_comparison_table(
        prs,
        "OSS reference stacks -- honest comparison",
        ["Dimension", "Frappe CRM (Py / Vue)", "Plane (TS)", "OpenProject (Rails)"],
        [
            ["Domain fit", "Sales CRM (leads, deals, calls, SLAs)", "Project / issue tracking", "Project / portfolio / Gantt"],
            ["Language / stack", "Python + Vue 3 + MariaDB / PG", "Next.js + Django + PG", "Rails + Angular + PG"],
            ["Licence", "AGPL-3", "AGPL-3 (+ Enterprise tier)", "GPL-3 (+ Enterprise add-ons)"],
            ["Extensibility", "Metadata-driven DocTypes -- no fork", "Code forks; plugin API young", "Ruby plugins; Rails skills scarce"],
            ["Multi-tenant", "Sites-per-tenant, each own DB", "Workspaces inside one DB", "Typically single-tenant"],
            ["Integrations", "Twilio, WhatsApp, ERPNext, IMAP/SMTP, webhooks, RPC", "REST + webhooks", "REST + webhooks; SAML/LDAP"],
            ["CBRE fit", "Closest infra match -- CRM + extensible", "No -- not a CRM", "Good for projects, heavy to extend"],
        ],
    )

    # 7. Target architecture (full diagram)
    target_png = DIAGRAMS / "target_cbre.png"
    if target_png.exists():
        add_image_slide(
            prs,
            "Target architecture: OSS core + semantic layer + AI agents",
            target_png,
            caption="BFF + event bus + semantic layer (knowledge graph + CRE ontology). OSS core is the durable asset. AI agents reason over the ontology.",
        )
    else:
        add_bullet_slide(prs, "Target architecture (diagram missing)",
                         [f"Expected PNG at {target_png}"])

    # 8. Why this shape
    add_bullet_slide(
        prs,
        "Why this shape specifically",
        [
            "OSS core (Frappe CRM) is the system of record for accounts, projects, workflows -- a Postgres schema you own.",
            "Salesforce stays, but only for sales pipeline. SF becomes a consumer of accounts, not the master.",
            "A thin BFF + event bus is the integration surface. Swap or upgrade the OSS core later without touching clients.",
            "Tableau becomes read-only, fed from the event bus / warehouse. It stops being the front door.",
            "Learn from best-in-class CRE SaaS (e.g. Enaia) for domain UX, workflows, data models -- build those features on the OSS core you control.",
        ],
    )

    # 9. Semantic layer section header
    add_section_header(prs, "Semantic Layer",
                       "Knowledge graph + CRE ontology: the AI-readiness layer")

    # 10. Why a semantic layer
    add_bullet_slide(
        prs,
        "Why a semantic layer, not just a warehouse",
        [
            "Snowflake stores facts as rows. A knowledge graph stores facts AND their relationships.",
            '"Account A holds Lease B on Property C in Market D managed by Broker E" -- that is a graph, not a table.',
            "AI agents traverse context in a graph; they issue sequential SQL queries against a warehouse. Graph wins for reasoning.",
            "For a commercial real-estate business, the entity graph IS the business.",
            "The semantic layer sits between operational systems (CRM, SF, SharePoint) and AI -- single source of connected truth.",
        ],
        subtitle="The single most important structural decision for AI readiness.",
    )

    # 11. CRE ontology table
    add_comparison_table(
        prs,
        "CRE Ontology (OWL / SHACL): the domain vocabulary",
        ["Ontology class", "What it represents", "Key relationships"],
        [
            ["cre:Property", "Physical asset -- building, floor, unit", "hasLease, inMarket, managedBy"],
            ["cre:Lease", "Contractual tenancy", "hasTenant, onProperty, hasTerm"],
            ["cre:Tenant", "Occupier entity", "holdsLease, linkedAccount"],
            ["cre:Market", "Geographic / sector market", "containsProperty, hasComps"],
            ["cre:BrokerEngagement", "Advisory mandate", "forAccount, onProperty, assignedBroker"],
            ["cre:Portfolio", "Grouped set of properties", "containsProperty, ownedBy"],
            ["cre:Comp", "Comparable transaction", "inMarket, forPropertyType, atPrice"],
        ],
    )

    # 12. Entity resolver
    add_bullet_slide(
        prs,
        "Entity resolver: golden-record matching",
        [
            "Accounts live in CRM, Salesforce, SharePoint, and Snowflake -- often under different names.",
            '"CBRE Global Investors", "CBRE GI", and Salesforce Account #00341 resolve to one node in the graph.',
            "Fuzzy matching + deterministic rules (tax ID, DUNS, domain) across all sources.",
            "CDC events from CRM and SF feed the resolver continuously -- the graph stays current without batch ETL.",
            "This is the prerequisite for any AI agent that needs to answer questions spanning multiple systems.",
        ],
    )

    # 13. AI agents section
    add_section_header(prs, "AI Agents",
                       "Autonomous reasoning over the knowledge graph")

    # 14. AI agents + ontology
    add_bullet_slide(
        prs,
        "AI agents: what the ontology enables",
        [
            'Research: "Pull all lease expirations in the next 18 months for Portfolio X, find comps, draft a renewal strategy." Graph traversal, not API calls.',
            'Underwriting: "Given this property and market comps, score the deal." Agent reads the ontology to know what a comp is.',
            "Comp analysis: find the five most relevant comparables for a lease. Traverses cre:Comp -> cre:Market -> cre:PropertyType.",
            "Tool registry (MCP) gives agents structured write-back into CRM, SF, and the KG -- ontology-typed actions, not raw SQL.",
            "When new AI capabilities arrive (planning, multi-step reasoning), the ontology is already there. Upgrade models, not data structures.",
        ],
        subtitle="The ontology is the stable contract between your data and whatever AI runs on top of it.",
    )

    # 15. Sequencing
    add_roadmap_slide(prs)

    # 16. What NOT to do
    add_bullet_slide(
        prs,
        "What NOT to do",
        [
            "Do not let Salesforce lead a greenfield rebuild of account + project management.",
            "Do not pick 'fully custom composable platform' on day one -- use Frappe CRM for the 80%, learn CRE domain features from specialist SaaS like Enaia for the 20%.",
            "Do not treat Tableau as the action layer. Tableau is a read model.",
            "Do not skip the knowledge graph -- without the ontology, you get summarisation, not reasoning.",
            "Do not skip the event bus -- it is the piece that makes every later swap cheap.",
        ],
    )

    # 17. Reference runtime shape
    ref_png = DIAGRAMS / "reference_oss_crm.png"
    if ref_png.exists():
        add_image_slide(
            prs,
            "Reference runtime shape: Python back-end, Vue front-end",
            ref_png,
            caption="Frappe-shaped: Vue 3 SPA + Python web workers + RQ queue + Redis + MariaDB/Postgres + per-site file store.",
        )

    # 18. Alternative architecture views
    for alt_name in ["arch_exec_simple", "arch_data_flow", "arch_ai_stack", "arch_integration"]:
        alt_png = DIAGRAMS / f"{alt_name}.png"
        if alt_png.exists():
            captions = {
                "arch_exec_simple": "Executive view: three layers, one sentence each. For senior stakeholders who need the shape, not the detail.",
                "arch_data_flow": "Data flow view: how information moves from source systems through the semantic layer to AI agents and dashboards.",
                "arch_ai_stack": "AI stack view: from raw data through ontology and knowledge graph to autonomous agent actions.",
                "arch_integration": "Integration view: what connects to what, with the event bus and BFF as the two seams.",
            }
            add_image_slide(prs, f"Architecture option: {alt_name.replace('arch_', '').replace('_', ' ').title()}",
                            alt_png, caption=captions.get(alt_name, ""))

    # 19. Summary
    add_bullet_slide(
        prs,
        "The decision in one paragraph",
        [
            "Keep Salesforce for pipeline.",
            "Move accounts and projects onto a Frappe CRM open-source core -- learning domain features from best-in-class CRE SaaS like Enaia.",
            "Put a BFF, event bus, and knowledge graph between systems. Layer a CRE ontology (OWL/SHACL) so AI agents reason over the full entity graph.",
            "React-Native Material UI on top. Tableau goes read-only. Proprietary CRM is sunset over 12 months.",
            "Exit cost at any future point is a Postgres export + a portable ontology -- not a re-platform.",
        ],
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path} ({out_path.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    build(OUT)
